"""
7/28 Meeting 簡報產生器
內容：第二階段（7/16-7/27）成果 — AI伸展計劃、馬達控制整合、Physio Agent個人化、知識庫擴充
風格：白底、中文標楷體、英文/數字 Times New Roman（使用者固定偏好，不得使用深色背景）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

CN_FONT = '標楷體'
EN_FONT = 'Times New Roman'

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
NAVY       = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT     = RGBColor(0x2E, 0x74, 0xB5)
GREY       = RGBColor(0x59, 0x59, 0x59)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run_font(run, size=18, bold=False, color=BLACK, font_ascii=EN_FONT, font_ea=CN_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_ascii
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = parse_xml(f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_ea}"/>')
        rPr.append(ea)
    else:
        ea.set('typeface', font_ea)


def set_white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_white_bg(slide)
    return slide


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box, tf


def add_title_bar(slide, text, size=30, color=NAVY):
    box, tf = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=True, color=color)
    # 底線色塊
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(12.1), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return box


def add_bullets(slide, items, left=Inches(0.8), top=Inches(1.5), width=Inches(11.7), height=Inches(5.3),
                 size=18, level_sizes=None, line_gap=Pt(10)):
    """items: list of (text, level) ；level 0 = 主要項目, level 1 = 次項目"""
    box, tf = add_textbox(slide, left, top, width, height)
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = line_gap
        bullet = '● ' if level == 0 else '– '
        run = p.add_run()
        run.text = bullet + text
        sz = size if level == 0 else max(size - 3, 14)
        set_run_font(run, size=sz, bold=(level == 0), color=BLACK if level == 0 else GREY)
    return box


def add_quote_box(slide, label, quote_text, top, color, height=Inches(1.5)):
    box, tf = add_textbox(slide, Inches(0.8), top, Inches(11.7), height)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    set_run_font(r1, size=18, bold=True, color=color)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = '「' + quote_text + '」'
    set_run_font(r2, size=15, bold=False, color=BLACK)

    # 左側色條
    bar = slide.shapes.add_shape(1, Inches(0.6), top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return box


def add_table(slide, headers, rows, top=Inches(1.6), left=Inches(1.5), width=Inches(10.3), row_h=Inches(0.55)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_h * n_rows
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        set_run_font(run, size=16, bold=True, color=WHITE)

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREY if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            set_run_font(run, size=15, bold=False, color=BLACK)
    return table


def add_footer(slide, text):
    box, tf = add_textbox(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    set_run_font(run, size=11, color=GREY)


# ── 建立簡報 ──────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# 第1頁：標題
s = add_slide(prs)
box, tf = add_textbox(slide=s, left=Inches(1), top=Inches(2.5), width=Inches(11.3), height=Inches(2.5),
                       anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = '智慧辦公椅專題 — 第二階段成果報告'
set_run_font(r, size=34, bold=True, color=NAVY)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(14)
r2 = p2.add_run()
r2.text = 'AI 伸展計劃 × 馬達控制整合 × Physio Agent 個人化'
set_run_font(r2, size=20, color=ACCENT)

p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)
r3 = p3.add_run()
r3.text = '2026/07/28  ｜  陳翊昕'
set_run_font(r3, size=16, color=GREY)

# 第2頁：本階段任務回顧
s = add_slide(prs)
add_title_bar(s, '本階段任務回顧（7/16 - 7/27）')
add_bullets(s, [
    ('AI 伸展計劃功能：查詢近 7 天紀錄 → 統計壞坐姿頻率 → 回傳動作 + 影片', 0),
    ('馬達控制整合：新增 motor_logs、椅背感測 MQTT 擴充、觸發 API、決策邏輯', 0),
    ('Physio Agent 個人化：帶入使用者身高體重 + 壞坐姿統計', 0),
    ('知識庫擴充：新增 10+ 份伸展／人體工學文獻，優化 RAG 檢索', 0),
])
add_footer(s, '2/11')

# 第3頁：功能一
s = add_slide(prs)
add_title_bar(s, '功能一：AI 個人化伸展計劃')
add_bullets(s, [
    ('GET /api/agent/stretch-plan', 0),
    ('查詢使用者近 7 天坐姿紀錄，排除正常坐姿與空椅', 1),
    ('統計各壞坐姿出現次數與佔比', 1),
    ('依最頻繁的壞坐姿，從 20+ 動作資料庫中挑出對應伸展動作（含 YouTube 連結）', 1),
])
add_footer(s, '3/11')

# 第4頁：實測結果 - 伸展計劃
s = add_slide(prs)
add_title_bar(s, '實測結果：伸展計劃統計')
add_bullets(s, [
    ('測試帳號近 7 天資料：前傾 4 次（57%）、久坐／後仰／左傾 各 1 次（14%）', 0),
    ('系統自動推薦 10 個對應動作，例如「頸部左右旋轉」「胸部擴展伸展」，皆附 YouTube 影片連結', 0),
], top=Inches(1.5), height=Inches(1.8))
add_table(
    s,
    headers=['坐姿類別', '出現次數', '佔比'],
    rows=[
        ['前傾（烏龜頸）', '4', '57%'],
        ['久坐未動', '1', '14%'],
        ['過度後仰', '1', '14%'],
        ['身體左傾', '1', '14%'],
    ],
    top=Inches(3.5), left=Inches(2.9), width=Inches(7.5), row_h=Inches(0.5),
)
add_footer(s, '4/11')

# 第5頁：功能二
s = add_slide(prs)
add_title_bar(s, '功能二：馬達控制整合')
add_bullets(s, [
    ('新增 MotorLog model，記錄每次觸發的坐姿、馬達清單、原因', 0),
    ('MQTT subscriber 擴充支援 smartchair/sensor/back（椅背感測資料）', 0),
    ('POST /api/motor/trigger：依坐姿類別自動選擇對應馬達', 0),
], top=Inches(1.5), height=Inches(1.9))
add_table(
    s,
    headers=['坐姿', '觸發馬達'],
    rows=[
        ['前傾', 'M1、M2（左右手軸）'],
        ['後仰', 'M3、M4（左右腰部）'],
        ['左傾', 'M2、M4（對側矯正）'],
        ['右傾', 'M1、M3（對側矯正）'],
        ['久坐', 'M1 ~ M4 全部'],
    ],
    top=Inches(3.4), left=Inches(2.9), width=Inches(7.5), row_h=Inches(0.5),
)
add_footer(s, '5/11')

# 第6頁：實測結果 - 馬達觸發
s = add_slide(prs)
add_title_bar(s, '實測結果：馬達觸發')
add_bullets(s, [
    ('輸入：{"posture": "forward"}', 0),
    ('輸出：{"motors": ["M1","M2"], "triggered": true,', 0),
    ('        "message": "馬達觸發：頭部前傾（M1、M2）"}', 1),
    ('資料庫 MotorLog 同步寫入紀錄，可供後續分析觸發頻率', 0),
], top=Inches(1.6), height=Inches(3))
add_footer(s, '6/11')

# 第7頁：功能三
s = add_slide(prs)
add_title_bar(s, '功能三：Physio Agent 個人化升級')
add_bullets(s, [
    ('由單次問答升級為 ReAct 迴圈（Thought → Action → Observation → Thought）', 0),
    ('新增身高體重（BMI）與近 7 天壞坐姿統計整合進 get_posture_history', 0),
    ('流程：讀取個人化資料 → 知識庫檢索 → 觸發馬達 → 再次查詢驗證改善（雙向回授）', 0),
], top=Inches(1.5), height=Inches(2.2))

steps = ['① 讀取個人化資料\n(BMI + 7天統計)', '② 知識庫檢索\n(RAG)', '③ 觸發馬達\n(震動提醒)', '④ 驗證改善\n(再次查詢)']
box_w, gap = Inches(2.5), Inches(0.5)
start_x = Inches(0.9)
for i, txt in enumerate(steps):
    x = start_x + i * (box_w + gap)
    shp = s.shapes.add_shape(1, x, Inches(4.3), box_w, Inches(1.3))
    shp.fill.solid()
    shp.fill.fore_color.rgb = LIGHT_GREY
    shp.line.color.rgb = ACCENT
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = txt
    set_run_font(run, size=14, bold=True, color=NAVY)
    if i < len(steps) - 1:
        arrow_box, atf = add_textbox(s, x + box_w, Inches(4.55), gap, Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = '→'
        set_run_font(ar, size=20, bold=True, color=ACCENT)
add_footer(s, '7/11')

# 第8頁：實測結果 - 個人化前後對比
s = add_slide(prs)
add_title_bar(s, '實測結果：個人化前後對比')
add_quote_box(s, '改版前：', '只回傳「最近 5 筆坐姿紀錄」，無體型資訊，建議千篇一律', top=Inches(1.55), color=GREY, height=Inches(1.1))
add_quote_box(
    s, '改版後（真實回應節錄）：',
    '根據您的身高體重（170cm, 65kg, BMI 22.5）與坐姿紀錄……您的近 7 天坐姿統計也顯示頭部前傾佔了 57%，是您最頻繁的不良坐姿',
    top=Inches(2.9), color=ACCENT, height=Inches(2.0),
)
add_bullets(s, [
    ('證明建議是即時運算出來的個人化結果，非套版文字', 0),
], top=Inches(5.3), height=Inches(0.8))
add_footer(s, '8/11')

# 第9頁：功能四
s = add_slide(prs)
add_title_bar(s, '功能四：知識庫擴充')
add_bullets(s, [
    ('新增 10 份文獻：桌椅人體工學、腕隧道症候群、下背痛、坐站交替、久坐中斷建議、', 0),
    ('辦公室瑜伽、核心穩定、胸椎呼吸、穿戴式震動回饋實證、椅子扶手腰靠調整', 1),
    ('全部引用真實學術／機構來源：Mayo Clinic、CCOHS、WHO Guidelines、PubMed、', 0),
    ('PMC 系統性回顧、MDPI 等', 1),
    ('順帶修復問題：發現本機 FAISS 向量索引停留在 5/14（比多次知識庫更新都舊），', 0),
    ('代表 Agent 一直在用過時資料回答，已重建索引修復', 1),
], top=Inches(1.5), height=Inches(4.5))
add_footer(s, '9/11')

# 第10頁：實測結果 - 新知識庫檢索
s = add_slide(prs)
add_title_bar(s, '實測結果：新知識庫檢索')
add_bullets(s, [
    ('問題：「我的椅子扶手跟腰靠應該怎麼調整比較好？」', 0),
], top=Inches(1.55), height=Inches(0.8))
add_quote_box(
    s, '回應精準引用新增的 chair_adjustment.txt：',
    '腰靠高度：站立時觸摸腰椎自然前凸位置……扶手高度：放鬆肩膀，調整至恰好托住手肘',
    top=Inches(2.5), color=ACCENT, height=Inches(1.8),
)
add_bullets(s, [
    ('證明新文獻確實被檢索與引用，非查無資料', 0),
], top=Inches(4.6), height=Inches(0.8))
add_footer(s, '10/11')

# 第11頁：Q&A
s = add_slide(prs)
box, tf = add_textbox(s, Inches(1), Inches(3), Inches(11.3), Inches(1.5), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'Q & A'
set_run_font(r, size=40, bold=True, color=NAVY)
add_footer(s, '11/11')

prs.save('report_0728.pptx')
print('已產生 report_0728.pptx')
