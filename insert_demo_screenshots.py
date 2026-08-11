"""
把 Postman 截圖嵌入 report_0811.pptx 的截圖佔位頁（第 8-10 頁）。

不會重新產生整份簡報，只在既有的 report_0811.pptx 上動刀，
所以先前若已在 PowerPoint 手動調整過其他頁面也不會被覆蓋。

執行方式（依序對應 STEP A / STEP B / STEP C 三張截圖）：
    python insert_demo_screenshots.py step_a.png step_b.png step_c.png
"""

import sys
from pptx import Presentation
from pptx.util import Inches

PPTX_PATH = 'report_0811.pptx'
# 對應 gen_0811_ppt.py 建立佔位框時的位置與大小
PLACEHOLDER_LEFT = Inches(1.4)
PLACEHOLDER_TOP = Inches(1.5)
PLACEHOLDER_WIDTH = Inches(10.5)
PLACEHOLDER_HEIGHT = Inches(5.2)
HINT_MARKER = '貼上 Postman 截圖'


def main():
    if len(sys.argv) != 4:
        print('用法：python insert_demo_screenshots.py step_a.png step_b.png step_c.png')
        sys.exit(1)

    image_paths = sys.argv[1:4]
    prs = Presentation(PPTX_PATH)

    target_slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and HINT_MARKER in shape.text_frame.text:
                target_slides.append((slide, shape))
                break

    if len(target_slides) != 3:
        print(f'[Error] 預期找到 3 個佔位框，實際找到 {len(target_slides)} 個。'
              f'可能已被手動編輯過，請改用 PowerPoint 手動貼圖。')
        sys.exit(1)

    for (slide, placeholder_shape), img_path in zip(target_slides, image_paths):
        left, top = placeholder_shape.left, placeholder_shape.top
        max_w, max_h = placeholder_shape.width, placeholder_shape.height

        sp = placeholder_shape._element
        sp.getparent().remove(sp)

        pic = slide.shapes.add_picture(img_path, left, top)
        scale = min(max_w / pic.width, max_h / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = left + (max_w - pic.width) // 2
        pic.top = top + (max_h - pic.height) // 2
        print(f'已插入 {img_path}')

    prs.save(PPTX_PATH)
    print(f'完成！已更新 {PPTX_PATH}')


if __name__ == '__main__':
    main()
